#!/usr/bin/env python3
"""
La-Z-Boy Branded PowerPoint Generator

Uses the official La-Z-Boy .potx template to create on-brand presentations.
Requires: pip install python-pptx

Usage:
    python create_pptx.py --output my_deck.pptx --slides slides.json

The slides.json format:
[
    {"layout": "title", "title": "Q1 Results", "subtitle": "Engineering Team"},
    {"layout": "divider", "title": "Agenda"},
    {"layout": "content", "title": "Key Updates", "body": "- Item 1\\n- Item 2\\n- Item 3"},
    {"layout": "two_column", "title": "Comparison", "left": "Column A text", "right": "Column B text"},
    {"layout": "image_text", "title": "Our Team", "body": "Description...", "image": "path/to/photo.jpg"},
    {"layout": "stats", "title": "KPIs", "body": "Revenue grew 15%...", "stat1": "$2.1B", "stat1_label": "Revenue"},
    {"layout": "end"}
]
"""

import argparse
import json
import os
import shutil
import sys
import tempfile
import zipfile
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


# Layout index mapping for La-Z-Boy brand (default)
LAYOUTS = {
    "title": 0,           # Title Slide — cover
    "title_alt1": 1,      # 1_Title Slide — alternate cover
    "title_alt2": 2,      # 2_Title Slide — alternate cover
    "divider": 3,         # Divider Slide 1 — section break
    "divider_alt": 4,     # Divider Slide 2 — section break alt
    "content_cream": 5,   # Title and Text Cream
    "content_white": 6,   # Title and Text White
    "content": 7,         # 1_Title and Text Cream (full height) — DEFAULT content
    "content_white_full": 8,  # 1_Title and Text White (full height)
    "two_column": 9,      # Two Column Text
    "round_image": 10,    # Text and Round Image
    "image_text_white": 11,   # Image and Text White
    "image_text": 12,     # Image and Text Cream
    "full_image": 13,     # Large Picture
    "profile": 14,        # Profile Slide
    "title_only": 15,     # Title and Footer Only
    "two_col_content": 16,    # Two Column Content (with headers)
    "three_col_content": 17,  # Three Column Content
    "four_col_content": 18,   # Four Column Content
    "quote": 19,          # Quote Slide
    "process": 20,        # Process Slide (5 steps)
    "stats": 21,          # Statistics Slide
    "footer_only": 22,    # Footer Only
    "blank": 23,          # Blank
    "logo": 24,           # Logo Lockup
    "brands": 25,         # Portfolio of Brands
    "end": 26,            # End Slide
}


def convert_potx_to_pptx(potx_path: str) -> str:
    """Convert .potx template to .pptx for python-pptx compatibility."""
    tmp = tempfile.mktemp(suffix=".pptx")
    with open(potx_path, "rb") as f:
        data = f.read()
    bio = BytesIO(data)
    with zipfile.ZipFile(bio, "r") as zin:
        with zipfile.ZipFile(tmp, "w") as zout:
            for item in zin.infolist():
                content = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    content = content.replace(
                        b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                        b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                    )
                zout.writestr(item, content)
    return tmp


def get_template_path() -> str:
    """Locate the .potx template relative to this script."""
    script_dir = Path(__file__).parent.parent
    potx = script_dir / "assets" / "lazboy-template.potx"
    if potx.exists():
        return str(potx)
    raise FileNotFoundError(
        f"Template not found at {potx}. "
        "Ensure assets/lazboy-template.potx exists in the skill directory."
    )


def create_presentation(slides_data: list[dict], output_path: str, brand: str = "lazboy") -> str:
    """Create a branded PowerPoint presentation from structured slide data.

    Args:
        slides_data: List of slide dicts with 'layout' key and content fields.
        output_path: Path to save the .pptx file.
        brand: Brand variant — 'lazboy' (default), 'england', or 'joybird'.

    Returns:
        Path to the created .pptx file.
    """
    potx_path = get_template_path()
    pptx_path = convert_potx_to_pptx(potx_path)

    try:
        prs = Presentation(pptx_path)

        # Remove all pre-existing sample slides from the template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

        for slide_data in slides_data:
            layout_key = slide_data.get("layout", "content")
            layout_idx = LAYOUTS.get(layout_key, LAYOUTS["content"])

            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            _populate_slide(slide, slide_data, layout_key)

        prs.save(output_path)
        return output_path
    finally:
        os.unlink(pptx_path)


def _populate_slide(slide, data: dict, layout_key: str) -> None:
    """Populate a slide's placeholders based on the layout type."""
    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    body = data.get("body", "")

    # Title placeholder (PH0 in most layouts)
    if title and 0 in slide.placeholders:
        slide.placeholders[0].text = title

    # Subtitle (PH1 — title slides only)
    if subtitle and 1 in slide.placeholders:
        slide.placeholders[1].text = subtitle

    # Main body content
    if layout_key in ("content", "content_cream", "content_white", "content_white_full"):
        if body and 1 in slide.placeholders:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, line in enumerate(body.split("\n")):
                if i == 0:
                    tf.paragraphs[0].text = line.lstrip("- ")
                else:
                    p = tf.add_paragraph()
                    p.text = line.lstrip("- ")

    # Two column
    elif layout_key == "two_column":
        left = data.get("left", "")
        right = data.get("right", "")
        if left and 1 in slide.placeholders:
            slide.placeholders[1].text = left
        if right and 2 in slide.placeholders:
            slide.placeholders[2].text = right

    # Image + text
    elif layout_key in ("image_text", "image_text_white"):
        image_path = data.get("image", "")
        if body and 2 in slide.placeholders:
            slide.placeholders[2].text = body
        if image_path and os.path.exists(image_path) and 10 in slide.placeholders:
            slide.placeholders[10].insert_picture(image_path)

    # Stats
    elif layout_key == "stats":
        if body and 17 in slide.placeholders:
            slide.placeholders[17].text = body
        stat_pairs = [
            (data.get("stat1"), data.get("stat1_label"), 18, 20),
            (data.get("stat2"), data.get("stat2_label"), 21, 23),
        ]
        for value, label, val_ph, lbl_ph in stat_pairs:
            if value and val_ph in slide.placeholders:
                slide.placeholders[val_ph].text = value
            if label and lbl_ph in slide.placeholders:
                slide.placeholders[lbl_ph].text = label

    # Quote
    elif layout_key == "quote":
        quote_text = data.get("quote", body)
        attribution = data.get("attribution", "")
        if quote_text and 10 in slide.placeholders:
            slide.placeholders[10].text = quote_text
        if attribution and 2 in slide.placeholders:
            slide.placeholders[2].text = attribution


def main():
    parser = argparse.ArgumentParser(description="Create La-Z-Boy branded PowerPoint")
    parser.add_argument("--output", "-o", required=True, help="Output .pptx file path")
    parser.add_argument("--slides", "-s", required=True, help="JSON file with slide data")
    parser.add_argument("--brand", default="lazboy", choices=["lazboy", "england", "joybird"])
    args = parser.parse_args()

    with open(args.slides) as f:
        slides_data = json.load(f)

    output = create_presentation(slides_data, args.output, args.brand)
    print(f"Created: {output} ({len(slides_data)} slides)")


if __name__ == "__main__":
    main()
